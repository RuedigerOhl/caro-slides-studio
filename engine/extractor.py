"""Extract structured content from a source PPTX as JSON-serializable data.

Output per slide:
  {
    "index": 1,
    "title": "...",                     # heuristic
    "texts": ["..."],                   # all non-title text blocks
    "bullets": ["..."],                 # split if a single block looks like a list
    "images": [
      {"path": "extracted/img1.png", "left": 0.5, "top": 1.0, "width": 4, "height": 3}
    ],
    "raw_layout_name": "Titel und Inhalt",
  }
"""
from __future__ import annotations
import json
import shutil
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def emu_to_inch(emu):
    if emu is None:
        return None
    return round(emu / 914400, 3)


def looks_like_bullet_list(text: str) -> bool:
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if len(lines) < 2:
        return False
    bulletish = sum(1 for l in lines if len(l) < 200)
    return bulletish >= len(lines) * 0.7


def extract_deck(source_pptx: str, image_dir: str | None = None) -> list[dict]:
    src = Path(source_pptx)
    if image_dir:
        img_dir = Path(image_dir)
        img_dir.mkdir(parents=True, exist_ok=True)
        # unzip just media
        with zipfile.ZipFile(src) as z:
            media_files = [n for n in z.namelist() if n.startswith('ppt/media/')]
            for m in media_files:
                dest = img_dir / Path(m).name
                with z.open(m) as src_f, open(dest, 'wb') as dst_f:
                    shutil.copyfileobj(src_f, dst_f)

    prs = Presentation(str(src))
    slides_out = []
    for si, slide in enumerate(prs.slides, start=1):
        title = ''
        texts = []
        images = []
        layout_name = slide.slide_layout.name if slide.slide_layout else ''
        # find title placeholder
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if shape == getattr(slide.shapes, 'title', None):
                continue
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and t != title:
                    texts.append(t)
            # collect picture shapes (shape_type 13)
            if shape.shape_type == 13:
                try:
                    img_info = {
                        'name': shape.name,
                        'left': emu_to_inch(shape.left),
                        'top': emu_to_inch(shape.top),
                        'width': emu_to_inch(shape.width),
                        'height': emu_to_inch(shape.height),
                    }
                    if image_dir:
                        # try to save image
                        try:
                            blob = shape.image.blob
                            ext = shape.image.ext or 'png'
                            fname = f'slide{si}_{shape.shape_id}.{ext}'
                            (Path(image_dir) / fname).write_bytes(blob)
                            img_info['path'] = str(Path(image_dir) / fname)
                        except Exception:
                            pass
                    images.append(img_info)
                except Exception:
                    pass

        # heuristic: if title empty but first text looks short, treat as title
        if not title and texts:
            first = texts[0]
            if len(first) < 120 and '\n' not in first:
                title = first
                texts = texts[1:]

        bullets = []
        paragraphs = []
        for t in texts:
            if looks_like_bullet_list(t):
                bullets.extend([l.strip() for l in t.splitlines() if l.strip()])
            else:
                paragraphs.append(t)

        slides_out.append({
            'index': si,
            'title': title,
            'paragraphs': paragraphs,
            'bullets': bullets,
            'images': images,
            'raw_layout_name': layout_name,
        })

    return slides_out


if __name__ == '__main__':
    import sys
    src = sys.argv[1]
    img_dir = sys.argv[2] if len(sys.argv) > 2 else None
    data = extract_deck(src, img_dir)
    print(json.dumps(data, indent=2, ensure_ascii=False))
