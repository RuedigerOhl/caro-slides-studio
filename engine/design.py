"""Design tokens + reusable slide builders for Caros Slide Studio master.

Single source of truth. All transformations and new-slide generation render
through these builders so style is identical everywhere.
"""
from pathlib import Path
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---- design tokens
RED = RGBColor(0x8F, 0x15, 0x26)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADLINE_FONT = 'Cambria'
BODY_FONT = 'Calibri'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Logo
LOGO_PATH = Path(__file__).resolve().parent.parent / 'master' / 'karrierecoach-logo.png'
LOGO_W = Inches(1.9)
LOGO_H = Inches(0.31)
LOGO_LEFT = SLIDE_W - LOGO_W - Inches(0.4)
LOGO_TOP = Inches(0.3)

# Layout grid
LEFT_MARGIN = Inches(0.6)
RIGHT_MARGIN = Inches(0.6)
HL_LEFT = LEFT_MARGIN
HL_TOP = Inches(0.45)
HL_WIDTH = SLIDE_W - HL_LEFT - LOGO_W - Inches(0.6)
HL_HEIGHT = Inches(0.7)
RULE_Y = Inches(1.15)
RULE_LEFT = LEFT_MARGIN
RULE_WIDTH = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
CONTENT_TOP = Inches(1.75)
CONTENT_HEIGHT = SLIDE_H - CONTENT_TOP - Inches(0.6)
CONTENT_LEFT = LEFT_MARGIN
CONTENT_WIDTH = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN


def set_run(run, *, font=BODY_FONT, size=18, bold=False, color=DARK, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_logo(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), LOGO_LEFT, LOGO_TOP, width=LOGO_W, height=LOGO_H)


def add_rule(slide, *, left=None, y=None, width=None, weight=1.5):
    line = slide.shapes.add_connector(
        1,
        left if left is not None else RULE_LEFT,
        y if y is not None else RULE_Y,
        (left if left is not None else RULE_LEFT) + (width if width is not None else RULE_WIDTH),
        y if y is not None else RULE_Y,
    )
    line.line.color.rgb = RED
    line.line.width = Pt(weight)
    return line


def add_headline(slide, text, *, with_rule=True):
    tb = slide.shapes.add_textbox(HL_LEFT, HL_TOP, HL_WIDTH, HL_HEIGHT)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, font=HEADLINE_FONT, size=30, bold=True, color=RED)
    if with_rule:
        add_rule(slide)
    return tb


def add_footer_pagenum(slide, n_text):
    tb = slide.shapes.add_textbox(LEFT_MARGIN, SLIDE_H - Inches(0.4), Inches(5), Inches(0.3))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Carolin Martin · Karrierecoach München'
    set_run(r, font=BODY_FONT, size=9, color=GRAY)
    tb = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4), Inches(0.8), Inches(0.3))
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = str(n_text)
    set_run(r, font=BODY_FONT, size=9, color=GRAY)


def add_image_placeholder(slide, left, top, width, height, label='Bild-Platzhalter'):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = LIGHT_GRAY
    r = rect.text_frame.paragraphs[0].add_run()
    r.text = label
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(r, font=BODY_FONT, size=14, color=GRAY)
    return rect


# ---- High-level slide builders ----

def render_cover(slide, *, title, subtitle='', presenter='Carolin Martin',
                 sub_lines=('Karrierecoach München',), date=''):
    add_logo(slide)
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.5), Inches(11), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run(r, font=HEADLINE_FONT, size=54, bold=True, color=RED)
    add_rule(slide, left=LEFT_MARGIN, y=Inches(4.2), width=Inches(5.8))
    if subtitle:
        tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.55), Inches(11), Inches(0.6))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = subtitle
        set_run(r, font=BODY_FONT, size=22, color=DARK)
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.7), Inches(11), Inches(1.2))
    lines = [presenter, *sub_lines]
    if date:
        lines.append(date)
    for i, line_text in enumerate(lines):
        para = tb.text_frame.add_paragraph() if i else tb.text_frame.paragraphs[0]
        r = para.add_run()
        r.text = line_text
        set_run(r, font=BODY_FONT, size=16, color=GRAY)


def render_section(slide, *, number, title, subtitle='', page=None):
    add_logo(slide)
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.4), Inches(3.5), Inches(2.8))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = str(number)
    set_run(r, font=HEADLINE_FONT, size=160, bold=True, color=RED)
    tb = slide.shapes.add_textbox(Inches(4.2), Inches(3.1), Inches(8.5), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run(r, font=HEADLINE_FONT, size=44, bold=True, color=DARK)
    add_rule(slide, left=Inches(4.25), y=Inches(4.65), width=Inches(3.25))
    if subtitle:
        tb = slide.shapes.add_textbox(Inches(4.2), Inches(4.85), Inches(8.5), Inches(0.6))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = subtitle
        set_run(r, font=BODY_FONT, size=18, color=GRAY, italic=True)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_bullets(slide, *, title, bullets, page=None):
    add_logo(slide)
    add_headline(slide, title)
    body = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    body.text_frame.word_wrap = True
    body.text_frame.margin_left = 0
    for i, bp in enumerate(bullets):
        para = body.text_frame.add_paragraph() if i else body.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        r = para.add_run()
        r.text = '•  '
        set_run(r, font=BODY_FONT, size=20, color=RED, bold=True)
        r = para.add_run()
        r.text = bp
        set_run(r, font=BODY_FONT, size=20, color=DARK)
        para.space_after = Pt(10)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_paragraph(slide, *, title, paragraphs, page=None):
    add_logo(slide)
    add_headline(slide, title)
    body = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    body.text_frame.word_wrap = True
    body.text_frame.margin_left = 0
    for i, para_text in enumerate(paragraphs):
        para = body.text_frame.add_paragraph() if i else body.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        r = para.add_run()
        r.text = para_text
        set_run(r, font=BODY_FONT, size=18, color=DARK)
        para.space_after = Pt(14)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_two_column(slide, *, title, left_title, left_body, right_title, right_body, page=None):
    add_logo(slide)
    add_headline(slide, title)
    col_w = (CONTENT_WIDTH - Inches(0.4)) / 2
    for ci, (htitle, btext) in enumerate([(left_title, left_body), (right_title, right_body)]):
        x = CONTENT_LEFT + (col_w + Inches(0.4)) * ci
        if htitle:
            tb = slide.shapes.add_textbox(x, CONTENT_TOP, col_w, Inches(0.5))
            r = tb.text_frame.paragraphs[0].add_run()
            r.text = htitle
            set_run(r, font=HEADLINE_FONT, size=22, bold=True, color=RED)
            body_top = CONTENT_TOP + Inches(0.6)
        else:
            body_top = CONTENT_TOP
        tb = slide.shapes.add_textbox(x, body_top, col_w, Inches(4.5))
        tb.text_frame.word_wrap = True
        if isinstance(btext, list):
            for i, bp in enumerate(btext):
                para = tb.text_frame.add_paragraph() if i else tb.text_frame.paragraphs[0]
                r = para.add_run()
                r.text = '•  '
                set_run(r, font=BODY_FONT, size=18, color=RED, bold=True)
                r = para.add_run()
                r.text = bp
                set_run(r, font=BODY_FONT, size=18, color=DARK)
                para.space_after = Pt(8)
        else:
            r = tb.text_frame.paragraphs[0].add_run()
            r.text = btext
            set_run(r, font=BODY_FONT, size=18, color=DARK)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_quote(slide, *, quote, attribution='', page=None):
    add_logo(slide)
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.4), Inches(2), Inches(2.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = '„'
    set_run(r, font=HEADLINE_FONT, size=200, bold=True, color=RED)
    quote_left = LEFT_MARGIN + Inches(1.4)
    tb = slide.shapes.add_textbox(quote_left, Inches(2.6), CONTENT_WIDTH - Inches(1.4), Inches(3))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = quote
    set_run(r, font=HEADLINE_FONT, size=32, italic=True, color=DARK)
    if attribution:
        tb = slide.shapes.add_textbox(quote_left, Inches(5.8), CONTENT_WIDTH - Inches(1.4), Inches(0.5))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = '— ' + attribution
        set_run(r, font=BODY_FONT, size=16, color=GRAY)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_image_text(slide, *, title, side, image_path=None, text_title='', text_body='', page=None):
    """side='left' or 'right' — which side the image sits on."""
    add_logo(slide)
    add_headline(slide, title)
    img_w = Inches(6.0)
    img_h = Inches(4.6)
    gap = Inches(0.4)
    text_width = CONTENT_WIDTH - img_w - gap
    if side == 'left':
        img_left, text_left = CONTENT_LEFT, CONTENT_LEFT + img_w + gap
    else:
        text_left, img_left = CONTENT_LEFT, CONTENT_LEFT + text_width + gap
    _place_image_or_placeholder(slide, image_path, img_left, CONTENT_TOP, img_w, img_h)
    tb = slide.shapes.add_textbox(text_left, CONTENT_TOP, text_width, img_h)
    tb.text_frame.word_wrap = True
    if text_title:
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text_title
        set_run(r, font=HEADLINE_FONT, size=22, bold=True, color=RED)
        p_body = tb.text_frame.add_paragraph()
        p_body.space_before = Pt(12)
    else:
        p_body = tb.text_frame.paragraphs[0]
    r = p_body.add_run()
    r.text = text_body
    set_run(r, font=BODY_FONT, size=16, color=DARK)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_image_center(slide, *, title, image_path=None, caption='', page=None):
    add_logo(slide)
    add_headline(slide, title)
    img_w = Inches(8.5)
    img_h = Inches(4.4)
    img_left = (SLIDE_W - img_w) / 2
    _place_image_or_placeholder(slide, image_path, img_left, CONTENT_TOP, img_w, img_h, label='Schaubild / Foto')
    if caption:
        cap = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP + img_h + Inches(0.2), CONTENT_WIDTH, Inches(0.5))
        cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = cap.text_frame.paragraphs[0].add_run()
        r.text = caption
        set_run(r, font=BODY_FONT, size=14, color=GRAY, italic=True)
    if page is not None:
        add_footer_pagenum(slide, page)


def render_closing(slide, *, lines=('Carolin Martin · Karrierecoach München',
                                    'mail@karrierecoach-muenchen.de',
                                    'www.karrierecoach-muenchen.de'),
                   message='Danke.'):
    add_logo(slide)
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.9), Inches(12), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = message
    set_run(r, font=HEADLINE_FONT, size=96, bold=True, color=RED)
    add_rule(slide, left=LEFT_MARGIN, y=Inches(4.55), width=Inches(3.9))
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.85), Inches(12), Inches(1.5))
    for i, line_text in enumerate(lines):
        para = tb.text_frame.add_paragraph() if i else tb.text_frame.paragraphs[0]
        r = para.add_run()
        r.text = line_text
        set_run(r, font=BODY_FONT, size=16, color=GRAY)


def _place_image_or_placeholder(slide, image_path, left, top, width, height, label='Bild-Platzhalter'):
    """Place image inside the box without distortion: fit + center, keep aspect ratio."""
    if image_path and Path(image_path).exists():
        try:
            # Read native size to compute aspect-correct fit
            try:
                from PIL import Image as PILImage
                with PILImage.open(image_path) as im:
                    iw, ih = im.size
            except Exception:
                # fallback: add as-is at original aspect
                slide.shapes.add_picture(str(image_path), left, top, width=width)
                return
            if iw <= 0 or ih <= 0:
                raise ValueError('bad image dimensions')
            box_ratio = width / height
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                # image is wider — fit by width
                new_w = width
                new_h = int(width / img_ratio)
            else:
                # image is taller — fit by height
                new_h = height
                new_w = int(height * img_ratio)
            new_left = left + (width - new_w) // 2
            new_top = top + (height - new_h) // 2
            slide.shapes.add_picture(str(image_path), new_left, new_top,
                                     width=new_w, height=new_h)
            return
        except Exception:
            pass
    add_image_placeholder(slide, left, top, width, height, label=label)


# ---- Top-level dispatcher ----

def render_slide(prs, spec):
    """spec: {'layout': 'cover'|'section'|'bullets'|'paragraph'|'two_column'|'quote'|'image_left'|'image_right'|'image_center'|'closing', ...}"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    layout = spec.get('layout', 'paragraph')
    if layout == 'cover':
        render_cover(slide, title=spec.get('title', ''), subtitle=spec.get('subtitle', ''),
                     presenter=spec.get('presenter', 'Carolin Martin'),
                     sub_lines=tuple(spec.get('sub_lines', ['Karrierecoach München'])),
                     date=spec.get('date', ''))
    elif layout == 'section':
        render_section(slide, number=spec.get('number', '01'), title=spec.get('title', ''),
                       subtitle=spec.get('subtitle', ''), page=spec.get('page'))
    elif layout == 'bullets':
        render_bullets(slide, title=spec.get('title', ''),
                       bullets=spec.get('bullets', []), page=spec.get('page'))
    elif layout == 'paragraph':
        render_paragraph(slide, title=spec.get('title', ''),
                         paragraphs=spec.get('paragraphs', []), page=spec.get('page'))
    elif layout == 'two_column':
        render_two_column(slide, title=spec.get('title', ''),
                          left_title=spec.get('left_title', ''),
                          left_body=spec.get('left_body', ''),
                          right_title=spec.get('right_title', ''),
                          right_body=spec.get('right_body', ''),
                          page=spec.get('page'))
    elif layout == 'quote':
        render_quote(slide, quote=spec.get('quote', ''),
                     attribution=spec.get('attribution', ''), page=spec.get('page'))
    elif layout in ('image_left', 'image_right'):
        render_image_text(slide, title=spec.get('title', ''),
                          side='left' if layout == 'image_left' else 'right',
                          image_path=spec.get('image_path'),
                          text_title=spec.get('text_title', ''),
                          text_body=spec.get('text_body', ''), page=spec.get('page'))
    elif layout == 'image_center':
        render_image_center(slide, title=spec.get('title', ''),
                            image_path=spec.get('image_path'),
                            caption=spec.get('caption', ''), page=spec.get('page'))
    elif layout == 'closing':
        render_closing(slide, message=spec.get('message', 'Danke.'),
                       lines=tuple(spec.get('lines', [])) or None)
    return slide


def new_presentation():
    from pptx import Presentation
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p
