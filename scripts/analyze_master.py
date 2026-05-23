#!/usr/bin/env python3
"""Analyze a PPTX master to extract design system info."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from collections import Counter

def emu_to_inch(emu):
    return round(emu / 914400, 2) if emu else None

def analyze(path):
    p = Presentation(path)
    print(f"\n{'='*70}\nFILE: {Path(path).name}\n{'='*70}")
    print(f"Slide size: {emu_to_inch(p.slide_width)} x {emu_to_inch(p.slide_height)} inch")
    print(f"Slide count: {len(p.slides)}")
    print(f"Slide masters: {len(p.slide_masters)}")
    for sm_i, sm in enumerate(p.slide_masters):
        print(f"\n--- Master {sm_i} — Layouts ({len(sm.slide_layouts)}) ---")
        for li, layout in enumerate(sm.slide_layouts):
            placeholders = [(ph.placeholder_format.idx, ph.placeholder_format.type, ph.name) for ph in layout.placeholders]
            print(f"  [{li}] {layout.name}: {len(placeholders)} placeholders")
            for idx, typ, name in placeholders:
                print(f"        - ph{idx} {typ} '{name}'")

    # Collect fonts and colors across slides
    fonts = Counter()
    colors = Counter()
    shape_types = Counter()
    images = []
    layout_usage = Counter()

    for slide_i, slide in enumerate(p.slides):
        layout_usage[slide.slide_layout.name] += 1
        for shape in slide.shapes:
            shape_types[shape.shape_type] += 1
            if shape.shape_type == 13:  # PICTURE
                images.append((slide_i+1, shape.name, emu_to_inch(shape.left), emu_to_inch(shape.top), emu_to_inch(shape.width), emu_to_inch(shape.height)))
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
                        try:
                            if run.font.color and run.font.color.rgb:
                                colors[str(run.font.color.rgb)] += 1
                        except Exception:
                            pass

    print(f"\n--- Layouts used in slides ---")
    for name, n in layout_usage.most_common():
        print(f"  {n:3d}x  {name}")

    print(f"\n--- Fonts used ---")
    for font, n in fonts.most_common(15):
        print(f"  {n:4d}x  {font}")

    print(f"\n--- Explicit text colors (hex) ---")
    for c, n in colors.most_common(15):
        print(f"  {n:4d}x  #{c}")

    print(f"\n--- Images on slides (first 10) ---")
    for entry in images[:10]:
        print(f"  slide {entry[0]:3d}: {entry[1]} @ ({entry[2]},{entry[3]}) size {entry[4]}x{entry[5]}")
    if len(images) > 10:
        print(f"  ... +{len(images)-10} more")

    # Theme analysis
    print(f"\n--- Theme colors (from first master) ---")
    try:
        theme = p.slide_masters[0].element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme')
        if theme is not None:
            ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            for child in theme:
                tag = child.tag.replace(ns, '')
                srgb = child.find(f'{ns}srgbClr')
                sys_clr = child.find(f'{ns}sysClr')
                if srgb is not None:
                    print(f"  {tag}: #{srgb.get('val')}")
                elif sys_clr is not None:
                    print(f"  {tag}: sys {sys_clr.get('val')} (lastClr=#{sys_clr.get('lastClr')})")
    except Exception as e:
        print(f"  (could not parse: {e})")

    # Font scheme
    print(f"\n--- Theme fonts ---")
    try:
        fonts_el = p.slide_masters[0].element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme')
        if fonts_el is not None:
            ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            for kind in ['majorFont', 'minorFont']:
                fel = fonts_el.find(f'{ns}{kind}')
                if fel is not None:
                    latin = fel.find(f'{ns}latin')
                    if latin is not None:
                        print(f"  {kind}: {latin.get('typeface')}")
    except Exception as e:
        print(f"  (could not parse: {e})")

if __name__ == '__main__':
    for arg in sys.argv[1:]:
        analyze(arg)
