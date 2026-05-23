#!/usr/bin/env python3
"""Build Caros Slide Studio Master v2 from scratch.

Design system:
  - Format: 16:9 (13.333 x 7.5 inch)
  - Primary red: #8F1526
  - Headline font: Cambria
  - Body font: Calibri
  - Logo: Karrierecoach München top-right

Layouts:
  1. Title (cover) — full-width red headline + subtitle + logo
  2. Section divider — large red number + section title on white
  3. Content — title with red underline + body content
  4. Two-column — title + two content blocks side by side
  5. Quote / callout — large centered quote in red
  6. Image + text — left image, right text
  7. Closing slide — thank-you + logo
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- design tokens
RED = RGBColor(0x8F, 0x15, 0x26)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADLINE_FONT = 'Cambria'
BODY_FONT = 'Calibri'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
LOGO_PATH = Path(__file__).resolve().parent.parent / 'master' / 'karrierecoach-logo.png'

# Logo placement: top-right, ~1.8" wide
LOGO_W = Inches(1.9)
LOGO_H = Inches(0.31)  # aspect ratio 6.12:1
LOGO_LEFT = SLIDE_W - LOGO_W - Inches(0.4)
LOGO_TOP = Inches(0.3)

# Unified left margin — headline + rule + content all share this x
LEFT_MARGIN = Inches(0.6)
RIGHT_MARGIN = Inches(0.6)

# Headline area (line higher, content gets breathing room below)
HL_LEFT = LEFT_MARGIN
HL_TOP = Inches(0.45)
HL_WIDTH = SLIDE_W - HL_LEFT - LOGO_W - Inches(0.6)
HL_HEIGHT = Inches(0.7)
RULE_Y = Inches(1.15)
RULE_LEFT = LEFT_MARGIN
RULE_WIDTH = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN

# Content area below rule — gap below line for breathing room
CONTENT_TOP = Inches(1.75)
CONTENT_HEIGHT = SLIDE_H - CONTENT_TOP - Inches(0.6)
CONTENT_LEFT = LEFT_MARGIN
CONTENT_WIDTH = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN


def set_run(run, *, font=BODY_FONT, size=18, bold=False, color=DARK, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_logo(slide):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), LOGO_LEFT, LOGO_TOP, width=LOGO_W, height=LOGO_H)


def add_accent_square(slide, x, y, size=Inches(0.18)):
    """Small red square — visual rhyme with the logo's red block."""
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    sq.fill.solid()
    sq.fill.fore_color.rgb = RED
    sq.line.fill.background()
    return sq


def add_headline(slide, text, *, with_rule=True):
    tb = slide.shapes.add_textbox(HL_LEFT, HL_TOP, HL_WIDTH, HL_HEIGHT)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, font=HEADLINE_FONT, size=30, bold=True, color=RED)
    if with_rule:
        add_rule(slide)
    return tb


def add_rule(slide):
    # Slim red rule, slight gradient feel: use a thin shape instead of connector
    line = slide.shapes.add_connector(1, RULE_LEFT, RULE_Y, RULE_LEFT + RULE_WIDTH, RULE_Y)
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)
    return line


def add_body_textbox(slide, left, top, width, height, *, placeholder='Text hier …'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = placeholder
    set_run(run, font=BODY_FONT, size=18, color=DARK)
    return tb


def add_footer_pagenum(slide, n_text='1'):
    # left: brand text, aligned with LEFT_MARGIN
    tb = slide.shapes.add_textbox(LEFT_MARGIN, SLIDE_H - Inches(0.4), Inches(5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = 'Carolin Martin · Karrierecoach München'
    set_run(r, font=BODY_FONT, size=9, color=GRAY)
    # right: page number
    tb = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4), Inches(0.8), Inches(0.3))
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = n_text
    set_run(r, font=BODY_FONT, size=9, color=GRAY)


# ---- layouts as example slides

def build():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    blank = p.slide_layouts[6]  # blank layout

    # 1) TITLE / COVER
    s = p.slides.add_slide(blank)
    add_logo(s)
    # big title
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(2.5), Inches(11), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r = p1.add_run()
    r.text = 'Titel der Präsentation'
    set_run(r, font=HEADLINE_FONT, size=54, bold=True, color=RED)
    # rule below title
    line = s.shapes.add_connector(1, LEFT_MARGIN, Inches(4.2), Inches(6.4), Inches(4.2))
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)
    # subtitle (gap below rule)
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(4.55), Inches(11), Inches(0.6))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Untertitel oder Anlass'
    set_run(r, font=BODY_FONT, size=22, color=DARK)
    # presenter info
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(5.7), Inches(11), Inches(1.2))
    tf = tb.text_frame
    for line_text in ['Carolin Martin', 'Karrierecoach München', 'TT.MM.JJJJ']:
        para = tf.add_paragraph() if line_text != 'Carolin Martin' else tf.paragraphs[0]
        r = para.add_run()
        r.text = line_text
        set_run(r, font=BODY_FONT, size=16, color=GRAY)

    # 2) SECTION DIVIDER
    s = p.slides.add_slide(blank)
    add_logo(s)
    # big red number
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(2.4), Inches(3.5), Inches(2.8))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = '01'
    set_run(r, font=HEADLINE_FONT, size=160, bold=True, color=RED)
    # section title
    tb = s.shapes.add_textbox(Inches(4.2), Inches(3.1), Inches(8.5), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Kapitel-Überschrift'
    set_run(r, font=HEADLINE_FONT, size=44, bold=True, color=DARK)
    # short rule
    line = s.shapes.add_connector(1, Inches(4.25), Inches(4.65), Inches(7.5), Inches(4.65))
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)
    # subtitle line
    tb = s.shapes.add_textbox(Inches(4.2), Inches(4.85), Inches(8.5), Inches(0.6))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Kurze Einleitung in das Kapitel'
    set_run(r, font=BODY_FONT, size=18, color=GRAY, italic=True)
    add_footer_pagenum(s, '2')

    # 3) CONTENT (title + bullets)
    s = p.slides.add_slide(blank)
    add_logo(s)
    add_headline(s, 'Inhalts-Slide mit Bulletpoints')
    body = s.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf = body.text_frame
    tf.word_wrap = True
    bullets = [
        'Erster zentraler Punkt — kurz und prägnant',
        'Zweiter Punkt mit etwas mehr Detail im Satzbau',
        'Dritter Punkt schließt den Gedankengang ab',
        'Vierter optionaler Punkt für mehr Tiefe',
    ]
    for i, bp in enumerate(bullets):
        para = tf.add_paragraph() if i else tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        # red dot prefix using bullet char
        r = para.add_run()
        r.text = '•  '
        set_run(r, font=BODY_FONT, size=20, color=RED, bold=True)
        r = para.add_run()
        r.text = bp
        set_run(r, font=BODY_FONT, size=20, color=DARK)
        para.space_after = Pt(10)
    add_footer_pagenum(s, '3')

    # 4) TWO-COLUMN
    s = p.slides.add_slide(blank)
    add_logo(s)
    add_headline(s, 'Zwei-Spalten-Layout für Vergleich')
    col_w = (CONTENT_WIDTH - Inches(0.4)) / 2
    for ci, (htitle, btext) in enumerate([
        ('Linke Spalte', 'Hier kommt der linke Inhalt rein — z.B. Vorher, Pro, oder Variante A.'),
        ('Rechte Spalte', 'Und hier der rechte Gegen-Inhalt — Nachher, Contra, oder Variante B.'),
    ]):
        x = CONTENT_LEFT + (col_w + Inches(0.4)) * ci
        # sub-heading
        tb = s.shapes.add_textbox(x, CONTENT_TOP, col_w, Inches(0.5))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = htitle
        set_run(r, font=HEADLINE_FONT, size=22, bold=True, color=RED)
        # body
        tb = s.shapes.add_textbox(x, CONTENT_TOP + Inches(0.6), col_w, Inches(3.5))
        tb.text_frame.word_wrap = True
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = btext
        set_run(r, font=BODY_FONT, size=18, color=DARK)
    add_footer_pagenum(s, '4')

    # 5) QUOTE / CALLOUT
    s = p.slides.add_slide(blank)
    add_logo(s)
    # giant red quote mark — sits at LEFT_MARGIN
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(1.4), Inches(2), Inches(2.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = '„'
    set_run(r, font=HEADLINE_FONT, size=200, bold=True, color=RED)
    # quote text — indented from the quote mark, aligned to a clean column
    quote_left = LEFT_MARGIN + Inches(1.4)
    tb = s.shapes.add_textbox(quote_left, Inches(2.6), CONTENT_WIDTH - Inches(1.4), Inches(3))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Hier steht ein zentrales Zitat oder ein Leitsatz — wirkt durch Whitespace und große Schrift.'
    set_run(r, font=HEADLINE_FONT, size=32, italic=True, color=DARK)
    # attribution
    tb = s.shapes.add_textbox(quote_left, Inches(5.8), CONTENT_WIDTH - Inches(1.4), Inches(0.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = '— Quelle oder Sprecher:in'
    set_run(r, font=BODY_FONT, size=16, color=GRAY)
    add_footer_pagenum(s, '5')

    def add_image_placeholder(slide, left, top, width, height, label='Bild-Platzhalter'):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = LIGHT_GRAY
        rect.line.color.rgb = LIGHT_GRAY
        r = rect.text_frame.paragraphs[0].add_run()
        r.text = label
        rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(r, font=BODY_FONT, size=14, color=GRAY)
        return rect

    # 6a) BILD LINKS, TEXT RECHTS
    s = p.slides.add_slide(blank)
    add_logo(s)
    add_headline(s, 'Bild links, Text rechts')
    img_w = Inches(6.0)
    img_h = Inches(4.6)
    gap = Inches(0.4)
    text_left = CONTENT_LEFT + img_w + gap
    text_width = CONTENT_WIDTH - img_w - gap
    add_image_placeholder(s, CONTENT_LEFT, CONTENT_TOP, img_w, img_h)
    tb = s.shapes.add_textbox(text_left, CONTENT_TOP, text_width, img_h)
    tb.text_frame.word_wrap = True
    p1 = tb.text_frame.paragraphs[0]
    r = p1.add_run()
    r.text = 'Bildunterschrift / Kontext'
    set_run(r, font=HEADLINE_FONT, size=22, bold=True, color=RED)
    p2 = tb.text_frame.add_paragraph()
    p2.space_before = Pt(12)
    r = p2.add_run()
    r.text = ('Hier kommt der erklärende Text zum Bild. Bullet-Punkte oder Fließtext '
              'sind beide möglich. Die Schrift bleibt ruhig, gut lesbar und einheitlich.')
    set_run(r, font=BODY_FONT, size=16, color=DARK)
    add_footer_pagenum(s, '6')

    # 6b) BILD RECHTS, TEXT LINKS — Text at LEFT_MARGIN (aligns with headline)
    s = p.slides.add_slide(blank)
    add_logo(s)
    add_headline(s, 'Bild rechts, Text links')
    text_width_l = CONTENT_WIDTH - img_w - gap
    img_left_r = CONTENT_LEFT + text_width_l + gap
    tb = s.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, text_width_l, img_h)
    tb.text_frame.word_wrap = True
    p1 = tb.text_frame.paragraphs[0]
    r = p1.add_run()
    r.text = 'Kernaussage zum Bild'
    set_run(r, font=HEADLINE_FONT, size=22, bold=True, color=RED)
    p2 = tb.text_frame.add_paragraph()
    p2.space_before = Pt(12)
    r = p2.add_run()
    r.text = ('Beschreibung links, Bild rechts. Gut für Storytelling: erst Argument lesen, '
              'dann visuell bestätigen. Funktioniert auch für Modelle, Schaubilder oder Portraits.')
    set_run(r, font=BODY_FONT, size=16, color=DARK)
    add_image_placeholder(s, img_left_r, CONTENT_TOP, img_w, img_h)
    add_footer_pagenum(s, '7')

    # 6c) BILD ZENTRIERT, GROSS
    s = p.slides.add_slide(blank)
    add_logo(s)
    add_headline(s, 'Bild zentriert')
    img_w_c = Inches(8.5)
    img_h_c = Inches(4.4)
    img_left_c = (SLIDE_W - img_w_c) / 2
    add_image_placeholder(s, img_left_c, CONTENT_TOP, img_w_c, img_h_c, 'Schaubild / Foto')
    # centered caption
    cap = s.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP + img_h_c + Inches(0.2), CONTENT_WIDTH, Inches(0.5))
    cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = cap.text_frame.paragraphs[0].add_run()
    r.text = 'Bildunterschrift: kurze Erklärung des Schaubilds'
    set_run(r, font=BODY_FONT, size=14, color=GRAY, italic=True)
    add_footer_pagenum(s, '8')

    # 7) CLOSING
    s = p.slides.add_slide(blank)
    add_logo(s)
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(2.9), Inches(12), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Danke.'
    set_run(r, font=HEADLINE_FONT, size=96, bold=True, color=RED)
    line = s.shapes.add_connector(1, LEFT_MARGIN, Inches(4.55), Inches(4.5), Inches(4.55))
    line.line.color.rgb = RED
    line.line.width = Pt(1.5)
    tb = s.shapes.add_textbox(LEFT_MARGIN, Inches(4.85), Inches(12), Inches(1.5))
    tf = tb.text_frame
    for i, line_text in enumerate([
        'Carolin Martin · Karrierecoach München',
        'mail@karrierecoach-muenchen.de',
        'www.karrierecoach-muenchen.de',
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run()
        r.text = line_text
        set_run(r, font=BODY_FONT, size=16, color=GRAY)

    out = Path(__file__).resolve().parent.parent / 'master' / 'caro-master-v2.pptx'
    p.save(str(out))
    print(f'wrote {out}')
    return out


if __name__ == '__main__':
    build()
